import streamlit as st
import openai
import os
from dotenv import load_dotenv
import spacy
import io
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

# Load environment variables and initialize API
load_dotenv()
api_key = os.getenv('OPENAI_API_KEY')
openai.api_key = api_key

# Load spaCy model for NLP tasks
nlp = spacy.load("en_core_web_sm")

# Functions to handle document reading and rule generation
def extract_text_from_file(uploaded_file):
    file_type = uploaded_file.name.split('.')[-1].lower()
    if file_type == 'pdf':
        return read_pdf(uploaded_file)
    elif file_type == 'docx':
        return read_docx(uploaded_file)
    elif file_type == 'pptx':
        return read_pptx(uploaded_file)
    else:
        return "Unsupported file type"

def read_pdf(file):
    with io.BytesIO(file.getvalue()) as f:
        reader = PdfReader(f)
        return ''.join([page.extract_text() for page in reader.pages if page.extract_text()])

def read_docx(file):
    doc = Document(file)
    return '\n'.join([para.text for para in doc.paragraphs if para.text])

def read_pptx(file):
    ppt = Presentation(file)
    return ' '.join([shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text")])

def generate_rules_from_doc(doc_content):
    doc = nlp(doc_content)
    return [f"Monitor for mention of {ent.text}" for ent in doc.ents if ent.label_ in ["ORG", "LAW", "PERSON"]]

# Updated risk scoring function to dynamically analyze AI response content
def calculate_risk_score(response, risk_keywords):
    risk_score = 0
    response_lower = response.lower()
    for keyword in risk_keywords:
        if keyword in response_lower:
            risk_score += 10  # Increment score for each risk-related keyword found
    return risk_score

# Define a set of risk-related keywords to look for in AI responses
risk_keywords = [
    "lack of alignment", "inadequate model", "lack of clarity", "failure to meet objectives",
    "deviation from guiding principles", "inadequate internal governance", "insufficient measures",
    "lack of stakeholder interaction", "inaccurate or outdated information", "non-compliance"
]

def get_chatgpt_response(prompt, rules, risk_threshold):
    context = f"The assistant must comply with the following rules: {', '.join(rules)} Adjust focus based on a risk threshold of {risk_threshold}."
    try:
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "system", "content": context}, {"role": "user", "content": prompt}],
        )
        return response.choices[0].message['content']
    except Exception as e:
        return f"An error occurred: {str(e)}"

# Streamlit UI
st.title('🛡️ Welcome to GuardianAI')
st.subheader('🤖 A Rule-Based Chat Prototype')

# Upload document and generate rules
uploaded_file = st.file_uploader("Upload a document with policies", type=['txt', 'pdf', 'docx', 'pptx'], key="file-uploader")
rules = []
if uploaded_file is not None:
    text = extract_text_from_file(uploaded_file)
    rules = generate_rules_from_doc(text)
    risk_threshold = st.slider("Set Risk Threshold", min_value=0, max_value=100, value=50, key="risk_slider")

    # Handle suggested questions and general user input
    response_section = st.empty()  # Placeholder for dynamic placement of user input and response

    # Display suggested questions after setting risk threshold
    st.subheader("Suggested Questions:")
    suggested_questions = ["What are the key compliance issues in this document?", "What are the potential risks in applying these policies?"]
    cols = st.columns(2)  # Create two columns
    for i, question in enumerate(suggested_questions):
        if cols[i].button(question):
            response = get_chatgpt_response(question, rules, risk_threshold)
            st.text_area("ChatGPT Response:", value=response, height=220)
            risk_score = calculate_risk_score(response, risk_keywords)
            st.metric("Risk Score for this question", risk_score)
            if risk_score >= risk_threshold:
                st.error("Risk Alert: This interaction has exceeded the risk threshold for this question!")
            else:
                st.success("Risk Level Acceptable for this question")
else:
    st.warning("Please upload a document to generate rules and set risk thresholds.")

# User input and response generation
user_input = st.text_input("Type your message to ChatGPT:", key="user_input")
if user_input:
    if rules:
        response = get_chatgpt_response(user_input, rules, risk_threshold)
        st.text_area("ChatGPT Response:", value=response, height=220)

        # Calculate and display risk
        risk_score = calculate_risk_score(response, risk_keywords)
        st.metric("Risk Score", risk_score)
        if risk_score >= risk_threshold:
            st.error("Risk Alert: This interaction has exceeded the risk threshold!")
        else:
            st.success("Risk Level Acceptable")
    else:
        st.error("Please upload a document to generate rules before interacting with ChatGPT.")

# Main sidebar setup
with st.sidebar:
    st.title("🛡️GuardianAI - Ensuring Gen AI")
    st.image("guardianai_logo.png")  # Ensure you have a logo image named 'guardianai_logo.png' in your working directory

    # User Guide Section
    st.title("User Guide")
    with st.expander("How to Use This AI Tool"):
        st.markdown("""
        **Welcome to GuardianAI!** Here’s how to get started:
        - Upload documents related to your organization's policies to analyze potential risks.
        - Use the sliders to set risk thresholds according to your compliance needs.
        - Interact with the AI chatbot by asking questions or providing scenarios to evaluate the AI's understanding and adherence to your policies.
        - The system will alert you if the content generated by the AI exceeds your set risk thresholds.
        """)

    # Admin Controls and Risk Management
    st.markdown("### Real-Time Monitoring and Admin Controls")
    st.write("Admins can review and approve monitoring rules, fine-tune alerting and prevention strategies based on risk thresholds, and manage user interactions in real time.")
    with st.expander("Project Overview"):
        st.write("""
            GuardianAI aims to create a comprehensive, real-time compliance and monitoring system designed for enterprises using Generative AI tools responsibly. The system offers risk-based alerting, preventative mechanisms, and automated monitoring rules that integrate seamlessly with an organization's existing IT infrastructure. This fosters a trustworthy AI environment by ensuring AI use complies with legal and ethical standards.
        """)

    # Feedback Section
    st.title("Feedback")
    rating = st.slider("Rate your experience", 1, 5, 4)
    feedback_text = st.text_input("Tell us more about your experience")
    if st.button("Submit Rating"):
        st.success(f"Thanks for rating us {rating} stars!")
        if feedback_text:
            st.write("Additional feedback:", feedback_text)
        st.markdown(
            "Do visit our [Website](https://github.com/MohamedFarhun/GuardianAI) for more information."
        )
